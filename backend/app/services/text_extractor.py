"""
Text Extraction Service
Handles extracting text content from various file formats.
Supports: PDF, DOCX, TXT, EML, XLSX, PPTX, images (via OCR)
"""

import os
import re
import email
import logging
from pathlib import Path
from typing import Tuple, Optional

logger = logging.getLogger(__name__)


class TextExtractor:
    """Extracts text from various document formats."""

    SUPPORTED_TYPES = {
        "pdf": "_extract_pdf",
        "docx": "_extract_docx",
        "doc": "_extract_docx",
        "txt": "_extract_txt",
        "eml": "_extract_email",
        "msg": "_extract_email",
        "xlsx": "_extract_excel",
        "xls": "_extract_excel",
        "pptx": "_extract_pptx",
        "csv": "_extract_txt",
        "json": "_extract_txt",
        "html": "_extract_html",
        "htm": "_extract_html",
        "rtf": "_extract_txt",
        "jpg": "_extract_image_ocr",
        "jpeg": "_extract_image_ocr",
        "png": "_extract_image_ocr",
        "tiff": "_extract_image_ocr",
        "bmp": "_extract_image_ocr",
    }

    def extract(self, file_path: str, file_type: str) -> Tuple[str, dict]:
        """
        Extract text and metadata from a file.
        Returns: (extracted_text, metadata_dict)
        """
        file_type = file_type.lower().strip(".")
        extractor_name = self.SUPPORTED_TYPES.get(file_type)

        if not extractor_name:
            raise ValueError(f"Unsupported file type: {file_type}")

        extractor = getattr(self, extractor_name)
        try:
            text, metadata = extractor(file_path)
            # Clean up extracted text
            text = self._clean_text(text)
            metadata["text_length"] = len(text)
            return text, metadata
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise

    def _extract_pdf(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from PDF files, with OCR fallback for scanned PDFs."""
        import pdfplumber

        text_parts = []
        metadata = {
            "page_count": 0,
            "has_images": False,
            "ocr_applied": False
        }

        try:
            with pdfplumber.open(file_path) as pdf:
                metadata["page_count"] = len(pdf.pages)

                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text() or ""
                    text_parts.append(page_text)

                    # Check if page has images
                    if page.images:
                        metadata["has_images"] = True

                # If very little text was extracted, try OCR
                full_text = "\n".join(text_parts)
                if len(full_text.strip()) < 50 and metadata["page_count"] > 0:
                    logger.info(f"Low text content in PDF, attempting OCR: {file_path}")
                    ocr_text = self._ocr_pdf(file_path)
                    if len(ocr_text) > len(full_text):
                        full_text = ocr_text
                        metadata["ocr_applied"] = True

            return full_text, metadata

        except Exception as e:
            logger.warning(f"pdfplumber failed for {file_path}: {e}, trying OCR")
            ocr_text = self._ocr_pdf(file_path)
            metadata["ocr_applied"] = True
            return ocr_text, metadata

    def _ocr_pdf(self, file_path: str) -> str:
        """OCR a PDF file using Tesseract."""
        try:
            import pytesseract
            from PIL import Image
            import subprocess
            import tempfile

            # Convert PDF to images using pdftoppm (if available) or Pillow
            text_parts = []

            try:
                # Try using pdf2image if available
                from pdf2image import convert_from_path
                images = convert_from_path(file_path, dpi=300)
                for img in images:
                    text = pytesseract.image_to_string(img)
                    text_parts.append(text)
            except ImportError:
                # Fallback: try converting first page with Pillow
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img)
                text_parts.append(text)

            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {e}")
            return ""

    def _extract_docx(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from DOCX files."""
        from docx import Document

        doc = Document(file_path)
        metadata = {
            "page_count": None,
            "has_images": False,
            "ocr_applied": False
        }

        text_parts = []
        for para in doc.paragraphs:
            text_parts.append(para.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_parts.append(cell.text)

        # Check for images
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                metadata["has_images"] = True
                break

        return "\n".join(text_parts), metadata

    def _extract_txt(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from plain text files."""
        metadata = {
            "page_count": None,
            "has_images": False,
            "ocr_applied": False
        }

        encodings = ["utf-8", "latin-1", "cp1252", "ascii"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                return text, metadata
            except (UnicodeDecodeError, UnicodeError):
                continue

        # Binary fallback
        with open(file_path, "rb") as f:
            text = f.read().decode("utf-8", errors="replace")
        return text, metadata

    def _extract_email(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from email files (.eml)."""
        metadata = {
            "page_count": None,
            "has_images": False,
            "ocr_applied": False,
            "email_from": None,
            "email_to": None,
            "email_subject": None,
            "email_date": None,
        }

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            msg = email.message_from_file(f)

        metadata["email_from"] = msg.get("From", "")
        metadata["email_to"] = msg.get("To", "")
        metadata["email_subject"] = msg.get("Subject", "")
        metadata["email_date"] = msg.get("Date", "")

        text_parts = []
        text_parts.append(f"From: {metadata['email_from']}")
        text_parts.append(f"To: {metadata['email_to']}")
        text_parts.append(f"Subject: {metadata['email_subject']}")
        text_parts.append(f"Date: {metadata['email_date']}")
        text_parts.append("---")

        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    payload = part.get_payload(decode=True)
                    if payload:
                        text_parts.append(payload.decode("utf-8", errors="replace"))
                elif content_type == "text/html":
                    payload = part.get_payload(decode=True)
                    if payload:
                        text_parts.append(self._strip_html(payload.decode("utf-8", errors="replace")))
        else:
            payload = msg.get_payload(decode=True)
            if payload:
                text_parts.append(payload.decode("utf-8", errors="replace"))

        return "\n".join(text_parts), metadata

    def _extract_excel(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from Excel files."""
        from openpyxl import load_workbook

        metadata = {
            "page_count": None,
            "has_images": False,
            "ocr_applied": False,
            "sheet_count": 0
        }

        wb = load_workbook(file_path, read_only=True, data_only=True)
        metadata["sheet_count"] = len(wb.sheetnames)

        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            text_parts.append(f"--- Sheet: {sheet_name} ---")
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                text_parts.append(row_text)

        wb.close()
        return "\n".join(text_parts), metadata

    def _extract_pptx(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from PowerPoint files."""
        from pptx import Presentation

        prs = Presentation(file_path)
        metadata = {
            "page_count": len(prs.slides),
            "has_images": False,
            "ocr_applied": False
        }

        text_parts = []
        for i, slide in enumerate(prs.slides):
            text_parts.append(f"--- Slide {i + 1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
                if shape.shape_type == 13:  # Picture
                    metadata["has_images"] = True

        return "\n".join(text_parts), metadata

    def _extract_html(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from HTML files."""
        metadata = {
            "page_count": None,
            "has_images": False,
            "ocr_applied": False
        }

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            html_content = f.read()

        text = self._strip_html(html_content)
        return text, metadata

    def _extract_image_ocr(self, file_path: str) -> Tuple[str, dict]:
        """Extract text from images using OCR."""
        import pytesseract
        from PIL import Image

        metadata = {
            "page_count": 1,
            "has_images": True,
            "ocr_applied": True
        }

        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        return text, metadata

    @staticmethod
    def _strip_html(html: str) -> str:
        """Remove HTML tags from text."""
        clean = re.sub(r"<[^>]+>", " ", html)
        clean = re.sub(r"\s+", " ", clean)
        return clean.strip()

    @staticmethod
    def _clean_text(text: str) -> str:
        """Clean extracted text by normalizing whitespace."""
        if not text:
            return ""
        # Remove null bytes
        text = text.replace("\x00", "")
        # Normalize line endings
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        # Remove excessive blank lines (more than 2)
        text = re.sub(r"\n{3,}", "\n\n", text)
        # Remove excessive spaces
        text = re.sub(r"[ \t]{3,}", "  ", text)
        return text.strip()


# Singleton instance
text_extractor = TextExtractor()
